from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)      # Dark navy
ACCENT = RGBColor(0x2E, 0x86, 0xAB)        # Teal blue
LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)      # Light gray-blue
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
SUBTITLE_CLR = RGBColor(0x5A, 0x5A, 0x5A)
ORANGE = RGBColor(0xE8, 0x6C, 0x00)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xC0, 0x39, 0x2B)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK_TEXT, bullet_char="•"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet_char}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(6)
    return txBox

def add_slide_number(slide, num, total=12):
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.4),
                f"{num} / {total}", font_size=10, color=SUBTITLE_CLR, alignment=PP_ALIGN.RIGHT)

# ========== SLIDE 1 — Cover ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, PRIMARY)
# Accent bar
add_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)
# Top accent line
add_rect(slide, Inches(0.15), Inches(2.8), Inches(3), Inches(0.06), ACCENT)

add_textbox(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.2),
            "Development of a Web Application\nfor Textile Workshop Management",
            font_size=36, bold=True, color=WHITE)

add_textbox(slide, Inches(1), Inches(3.2), Inches(6), Inches(2.5),
            "Student: Raiymbek Daniiar Uulu\nProgram: Computer Science\nGroup: COM22",
            font_size=20, color=RGBColor(0xBB, 0xCC, 0xDD))

add_textbox(slide, Inches(1), Inches(5.5), Inches(6), Inches(1.2),
            "Supervisor: [Supervisor Full Name]",
            font_size=18, color=RGBColor(0x99, 0xAA, 0xBB))

add_textbox(slide, Inches(1), Inches(6.3), Inches(6), Inches(0.8),
            "University: [University Name]  |  Year: 2026",
            font_size=16, color=RGBColor(0x88, 0x99, 0xAA))

# ========== SLIDE 2 — Introduction ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
            "Introduction", font_size=32, bold=True, color=WHITE)
add_slide_number(slide, 2)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.0),
            "The textile industry plays an important role in many local economies. Small and medium textile workshops produce garments such as T-shirts, dresses, uniforms, and other products.",
            font_size=18, color=DARK_TEXT)

add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.6),
            "However, many workshops still manage their operations manually using notebooks or spreadsheets. This leads to several problems:",
            font_size=18, color=DARK_TEXT)

add_bullet_list(slide, Inches(1.2), Inches(3.4), Inches(10), Inches(2.5), [
    "Poor order tracking",
    "Inaccurate inventory management",
    "Manual salary calculations",
    "Lack of production analytics"
], font_size=18, color=RED)

# Goal box
goal_shape = add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.2), LIGHT_BG)
add_textbox(slide, Inches(1.1), Inches(5.7), Inches(11), Inches(1.0),
            "Goal: Develop a web application that automates these processes and improves workshop management.",
            font_size=20, bold=True, color=PRIMARY)

# ========== SLIDE 3 — Problem Statement ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
            "Problem Statement", font_size=32, bold=True, color=WHITE)
add_slide_number(slide, 3)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
            "Many textile workshops face operational challenges:", font_size=20, bold=True, color=DARK_TEXT)

problems = [
    "Orders are tracked manually",
    "Materials are not properly monitored",
    "Production reports are recorded on paper",
    "Salary calculations are performed manually",
    "Management lacks clear analytical insights"
]
add_bullet_list(slide, Inches(1.2), Inches(2.4), Inches(10), Inches(3.0), problems, font_size=18)

add_rect(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.06), ACCENT)

add_textbox(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.6),
            "These problems reduce efficiency and increase the risk of errors.", font_size=18, color=SUBTITLE_CLR)

goal2 = add_rect(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.0), LIGHT_BG)
add_textbox(slide, Inches(1.1), Inches(6.1), Inches(11), Inches(0.8),
            "A digital management system can significantly improve production control and business transparency.",
            font_size=19, bold=True, color=PRIMARY)

# ========== SLIDE 4 — Research and Existing Solutions ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
            "Research and Existing Solutions", font_size=32, bold=True, color=WHITE)
add_slide_number(slide, 4)

# ERP Box
add_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.5), LIGHT_BG)
add_textbox(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.5),
            "ERP Systems", font_size=22, bold=True, color=PRIMARY)
add_textbox(slide, Inches(1.1), Inches(2.3), Inches(5), Inches(0.5),
            "SAP  •  Oracle ERP  •  Microsoft Dynamics", font_size=16, color=SUBTITLE_CLR)
add_textbox(slide, Inches(1.1), Inches(3.0), Inches(5), Inches(0.4),
            "Advantages:", font_size=16, bold=True, color=GREEN)
add_textbox(slide, Inches(1.1), Inches(3.4), Inches(5), Inches(0.4),
            "   Comprehensive business management", font_size=15, color=DARK_TEXT)
add_textbox(slide, Inches(1.1), Inches(3.9), Inches(5), Inches(0.4),
            "Disadvantages:", font_size=16, bold=True, color=RED)
add_bullet_list(slide, Inches(1.3), Inches(4.3), Inches(4.8), Inches(1.5), [
    "Complex implementation",
    "Expensive",
    "Not suitable for small workshops"
], font_size=15, color=DARK_TEXT, bullet_char="✗")

# MES Box
add_rect(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(4.5), LIGHT_BG)
add_textbox(slide, Inches(7.3), Inches(1.7), Inches(5), Inches(0.5),
            "Manufacturing Execution Systems (MES)", font_size=20, bold=True, color=PRIMARY)
add_textbox(slide, Inches(7.3), Inches(3.0), Inches(5), Inches(0.4),
            "Advantages:", font_size=16, bold=True, color=GREEN)
add_textbox(slide, Inches(7.3), Inches(3.4), Inches(5), Inches(0.4),
            "   Real-time production monitoring", font_size=15, color=DARK_TEXT)
add_textbox(slide, Inches(7.3), Inches(3.9), Inches(5), Inches(0.4),
            "Disadvantages:", font_size=16, bold=True, color=RED)
add_bullet_list(slide, Inches(7.5), Inches(4.3), Inches(4.8), Inches(1.5), [
    "Designed for large factories",
    "Difficult to adapt for small textile businesses"
], font_size=15, color=DARK_TEXT, bullet_char="✗")

# Conclusion
conc = add_rect(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.8), ACCENT)
add_textbox(slide, Inches(1.1), Inches(6.45), Inches(11), Inches(0.7),
            "Therefore, a simpler and more specialized solution is needed.",
            font_size=20, bold=True, color=WHITE)

# ========== SLIDE 5 — Proposed Solution ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
            "Proposed Solution", font_size=32, bold=True, color=WHITE)
add_slide_number(slide, 5)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
            "This project proposes a Web Application for Textile Workshop Management.",
            font_size=22, bold=True, color=PRIMARY)

add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(0.5),
            "The system will provide a centralized platform for managing:", font_size=18, color=DARK_TEXT)

# Feature cards
features = ["Orders", "Inventory", "Production", "Employees", "Financial Data"]
colors_feat = [ACCENT, GREEN, ORANGE, PRIMARY, RGBColor(0x8E, 0x44, 0xAD)]
for i, (feat, clr) in enumerate(zip(features, colors_feat)):
    x = Inches(0.8 + i * 2.4)
    add_rect(slide, x, Inches(3.3), Inches(2.1), Inches(1.8), clr)
    add_textbox(slide, x, Inches(3.8), Inches(2.1), Inches(0.8),
                feat, font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), LIGHT_BG)
add_textbox(slide, Inches(1.1), Inches(5.9), Inches(11), Inches(0.8),
            "The application will help workshop administrators track operations in real time and improve decision-making.",
            font_size=18, color=PRIMARY)

# ========== SLIDE 6 — Main Features ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
            "Main Features of the System", font_size=32, bold=True, color=WHITE)
add_slide_number(slide, 6)

modules = [
    ("Authentication", ["Login and password", "Password recovery", "Two-factor authentication"]),
    ("Dashboard", ["Monthly revenue", "Active orders", "Completed orders", "Defective items", "Inventory overview"]),
    ("Clients", ["Client database", "Order history"]),
    ("Orders", ["Order creation", "Product types", "Auto price calculation", "Order status tracking"]),
]

positions = [(0.8, 1.6), (4.0, 1.6), (7.2, 1.6), (10.4, 1.6)]
for (title, items), (x, y) in zip(modules, positions):
    add_rect(slide, Inches(x), Inches(y), Inches(2.9), Inches(5.2), LIGHT_BG)
    add_rect(slide, Inches(x), Inches(y), Inches(2.9), Inches(0.7), ACCENT)
    add_textbox(slide, Inches(x), Inches(y + 0.1), Inches(2.9), Inches(0.5),
                title, font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide, Inches(x + 0.2), Inches(y + 0.9), Inches(2.5), Inches(4.0),
                    items, font_size=14, color=DARK_TEXT, bullet_char="▸")

# ========== SLIDE 7 — Production and Inventory ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
            "Production and Inventory Management", font_size=32, bold=True, color=WHITE)
add_slide_number(slide, 7)

# Left: Inventory
add_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), LIGHT_BG)
add_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.7), GREEN)
add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5),
            "Inventory Management", font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(1.2), Inches(2.6), Inches(4.8), Inches(3.0), [
    "Material tracking",
    "Purchase operations",
    "Stock updates",
    "Usage history"
], font_size=18)

# Right: Technical Card
add_rect(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2), LIGHT_BG)
add_rect(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.7), ORANGE)
add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.5),
            "Technical Card (Bill of Materials)", font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(7.4), Inches(2.6), Inches(4.8), Inches(0.5),
            "Each product has a list of required materials.", font_size=16, color=DARK_TEXT)

add_textbox(slide, Inches(7.4), Inches(3.2), Inches(4.8), Inches(0.4),
            "Example — T-shirt production requires:", font_size=16, bold=True, color=PRIMARY)

add_bullet_list(slide, Inches(7.8), Inches(3.7), Inches(4.4), Inches(1.5), [
    "Fabric", "Thread", "Buttons"
], font_size=16, color=DARK_TEXT, bullet_char="▸")

add_rect(slide, Inches(7.2), Inches(5.2), Inches(5.0), Inches(1.2), ACCENT)
add_textbox(slide, Inches(7.4), Inches(5.3), Inches(4.6), Inches(1.0),
            "Materials are automatically deducted from inventory during production.",
            font_size=16, bold=True, color=WHITE)

# ========== SLIDE 8 — Employee and Salary ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
            "Employee and Salary Management", font_size=32, bold=True, color=WHITE)
add_slide_number(slide, 8)

# Left
add_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0), LIGHT_BG)
add_textbox(slide, Inches(1.1), Inches(1.8), Inches(5), Inches(0.5),
            "Employee Data", font_size=22, bold=True, color=PRIMARY)
add_bullet_list(slide, Inches(1.3), Inches(2.5), Inches(4.8), Inches(3.5), [
    "Name",
    "Position",
    "Phone number",
    "Piecework rate",
    "Hiring date"
], font_size=18)

# Right
add_rect(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.0), LIGHT_BG)
add_textbox(slide, Inches(7.3), Inches(1.8), Inches(5), Inches(0.5),
            "Salary Calculation", font_size=22, bold=True, color=PRIMARY)
add_textbox(slide, Inches(7.3), Inches(2.5), Inches(5), Inches(0.6),
            "Salary calculation is automated:", font_size=18, color=DARK_TEXT)

# Formula box
add_rect(slide, Inches(7.5), Inches(3.3), Inches(4.5), Inches(1.0), ACCENT)
add_textbox(slide, Inches(7.5), Inches(3.4), Inches(4.5), Inches(0.8),
            "Salary = produced items × piecework rate",
            font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(7.3), Inches(4.8), Inches(5), Inches(0.8),
            "The system generates monthly salary reports for employees.",
            font_size=18, color=DARK_TEXT)

# ========== SLIDE 9 — System Architecture ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
            "System Architecture", font_size=32, bold=True, color=WHITE)
add_slide_number(slide, 9)

# Architecture blocks
arch_items = [
    ("Frontend", "React", ACCENT),
    ("Backend", "Java Spring Boot", PRIMARY),
    ("Database", "PostgreSQL", GREEN),
    ("File Storage", "MinIO\n(S3-compatible)", ORANGE),
    ("Deployment", "Docker\ncontainers", RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (title, tech, clr) in enumerate(arch_items):
    x = Inches(0.8 + i * 2.5)
    # Box
    add_rect(slide, x, Inches(2.0), Inches(2.2), Inches(3.5), clr)
    add_textbox(slide, x, Inches(2.3), Inches(2.2), Inches(0.6),
                title, font_size=18, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.1) + x, Inches(3.2), Inches(2.0), Inches(1.8), WHITE)
    add_textbox(slide, Inches(0.1) + x, Inches(3.4), Inches(2.0), Inches(1.4),
                tech, font_size=17, bold=True, color=clr, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.9), LIGHT_BG)
add_textbox(slide, Inches(1.1), Inches(6.1), Inches(11), Inches(0.7),
            "This architecture ensures scalability, reliability, and easy deployment.",
            font_size=18, bold=True, color=PRIMARY)

# ========== SLIDE 10 — MVP ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
            "MVP (Minimum Viable Product)", font_size=32, bold=True, color=WHITE)
add_slide_number(slide, 10)

# Left: MVP
add_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.2), LIGHT_BG)
add_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.7), GREEN)
add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(0.5),
            "MVP Features", font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(1.2), Inches(2.6), Inches(4.8), Inches(4.0), [
    "Authentication system",
    "Dashboard",
    "Client management",
    "Order management",
    "Inventory tracking",
    "Production reporting"
], font_size=18, color=DARK_TEXT, bullet_char="✓")

# Right: Future
add_rect(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.2), LIGHT_BG)
add_rect(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.7), ACCENT)
add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.5),
            "Future Versions", font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_bullet_list(slide, Inches(7.4), Inches(2.6), Inches(4.8), Inches(4.0), [
    "Advanced analytics",
    "Mobile support",
    "Integration with accounting systems",
    "Multi-workshop management"
], font_size=18, color=DARK_TEXT, bullet_char="▸")

# ========== SLIDE 11 — Expected Results ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
            "Expected Results", font_size=32, bold=True, color=WHITE)
add_slide_number(slide, 11)

benefits = [
    ("Improved production control", GREEN),
    ("Accurate inventory tracking", ACCENT),
    ("Automated salary calculation", ORANGE),
    ("Better financial transparency", PRIMARY),
    ("Faster order processing", RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (text, clr) in enumerate(benefits):
    y = Inches(1.8 + i * 1.0)
    add_rect(slide, Inches(1.0), y, Inches(0.12), Inches(0.7), clr)
    add_rect(slide, Inches(1.3), y, Inches(10.5), Inches(0.7), LIGHT_BG)
    add_textbox(slide, Inches(1.6), y + Emu(Inches(0.1).emu), Inches(10), Inches(0.5),
                text, font_size=20, color=DARK_TEXT)

add_rect(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.9), PRIMARY)
add_textbox(slide, Inches(1.1), Inches(6.3), Inches(11), Inches(0.7),
            "The application will help textile workshop owners manage their operations more efficiently.",
            font_size=19, bold=True, color=WHITE)

# ========== SLIDE 12 — References ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.2), PRIMARY)
add_textbox(slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
            "References", font_size=32, bold=True, color=WHITE)
add_slide_number(slide, 12)

refs = [
    'Monk, E., & Wagner, B. (2013). Concepts in Enterprise Resource Planning. Cengage Learning.',
    'Laudon, K., & Laudon, J. (2020). Management Information Systems. Pearson.',
    'Kletti, J. (2007). Manufacturing Execution Systems. Springer.',
]

for i, ref in enumerate(refs):
    y = Inches(2.0 + i * 1.2)
    add_rect(slide, Inches(0.8), y, Inches(0.6), Inches(0.6), ACCENT)
    add_textbox(slide, Inches(0.8), y + Emu(Inches(0.05).emu), Inches(0.6), Inches(0.5),
                str(i + 1), font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.7), y + Emu(Inches(0.1).emu), Inches(10.5), Inches(0.5),
                ref, font_size=17, color=DARK_TEXT)

add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8),
            "Additional sources were analyzed to study manufacturing management systems and small business digitalization.",
            font_size=16, color=SUBTITLE_CLR)

# Thank you
add_rect(slide, Inches(3.5), Inches(6.2), Inches(6.3), Inches(0.9), PRIMARY)
add_textbox(slide, Inches(3.5), Inches(6.3), Inches(6.3), Inches(0.7),
            "Thank you for your attention!", font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Save
output_path = "/Users/raiymbekdaniiaruulu/IdealProjects/StartTups/claud/diplom/nuraiym/presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
